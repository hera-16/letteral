from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(31, 78, 121)

    subtitle_top = top + Inches(1.5)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = RGBColor(89, 89, 89)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

def add_content_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(9)
    title_height = Inches(0.8)

    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(31, 78, 121)

    shapes = slide.shapes
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(8.5)
    height = Inches(5)

    textbox = shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.font.color.rgb = RGBColor(64, 64, 64)

add_title_slide(prs, "Letteral", "組織問題の可視化プラットフォーム\n匿名で安全に、組織の問題を可視化")

add_content_slide(prs, "背景・課題", [
    "❌ 悩みが出づらい、本音が出ない",
    "",
    "❌ 進捗が他者に伝わらない",
    "",
    "❌ 組織規模に問わず、自分が所属しないグループでの問題が分からないことによる課題"
])

add_content_slide(prs, "解決コンセプト", [
    "💡 匿名性 × 階層的公開範囲制御",
    "",
    "   で\"安全な透明性\"を実現"
])

add_content_slide(prs, "主要機能（3本柱）", [
    "1️⃣ 匿名問題・進捗投稿",
    "   テンプレート / スレッド形式",
    "",
    "2️⃣ 階層・公開範囲・匿名ポリシー",
    "   任意深度の組織ツリー / 全社〜同期まで柔軟に制御",
    "",
    "3️⃣ 評価（自己評価）",
    "   投稿実績を評価材料として活用"
])

add_content_slide(prs, "誰に何が嬉しい？", [
    "👔 人事:",
    "   課題の早期可視化、離職リスク低減",
    "",
    "👨‍💼 管理職:",
    "   自己評価材料の多さ、状況把握",
    "",
    "👥 従業員:",
    "   匿名で本音、OKR紐付けで成長実感",
    "   他のグループの問題から自己対策"
])

add_content_slide(prs, "差別化（競合比較）", [
    "✅ 完全匿名性",
    "",
    "✅ 日次記録",
    "",
    "✅ 階層公開制御",
    "",
    "✅ 評価連動",
    "",
    "▶ これらを同時に満たすプラットフォーム"
])

add_content_slide(prs, "技術・信頼性", [
    "🔧 技術スタック:",
    "   Next.js / Spring Boot / JWT(RBAC)",
    "   WebSocket(STOMP) / MySQL / GitHub",
    "",
    "🔒 セキュリティ・非機能:",
    "   テナント分離(tenant_id)",
    "   監査ログ / モデレーション機能",
    "",
    "🎬 では、実際に",
    "   投稿→見える化→管理側の活用",
    "   を5分でお見せします"
])

prs.save('c:\\Users\\User\\OneDrive\\hera-16\\チャレキャラ\\Letteral_Presentation.pptx')
print("Presentation created: Letteral_Presentation.pptx")
